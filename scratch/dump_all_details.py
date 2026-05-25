import sys
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

def dump_all_shapes_and_text(pptx_path, out_file):
    prs = Presentation(pptx_path)
    slide_width_in = prs.slide_width / 914400
    mid_x_in = slide_width_in / 2
    
    out_file.write(f"Presentation: {pptx_path}\n")
    out_file.write(f"Slide dimensions: {slide_width_in:.2f} x {prs.slide_height / 914400:.2f} inches\n")
    out_file.write(f"Midpoint X: {mid_x_in:.2f} inches\n")
    out_file.write("=" * 60 + "\n")
    
    def process_shape(shape, slide_idx, indent=""):
        shape_type = shape.shape_type
        shape_name = shape.name
        
        # Calculate coordinates in inches
        try:
            left_in = shape.left / 914400
            top_in = shape.top / 914400
            width_in = shape.width / 914400
            height_in = shape.height / 914400
            pos_str = f"[{left_in:.2f}, {top_in:.2f}, {width_in:.2f}, {height_in:.2f}]"
            is_on_right = (left_in + width_in / 2) > mid_x_in
        except Exception:
            pos_str = "[unknown position]"
            is_on_right = False

        text_content = ""
        
        # 1. Standard text frame
        if shape.has_text_frame:
            text_content = shape.text_frame.text.strip()
            
        # 2. Table shape
        elif shape.has_table:
            table = shape.table
            table_text = []
            for r_idx, row in enumerate(table.rows):
                row_text = []
                for c_idx, cell in enumerate(row.cells):
                    row_text.append(cell.text.strip())
                table_text.append(f"Row {r_idx}: " + " | ".join(row_text))
            text_content = "\n".join(table_text)
            
        # 3. Group shape (recursive check)
        elif shape_type == MSO_SHAPE_TYPE.GROUP:
            out_file.write(f"{indent}Group Shape '{shape_name}' at {pos_str}:\n")
            for sub_shape in shape.shapes:
                process_shape(sub_shape, slide_idx, indent + "  ")
            return

        # 4. Chart
        elif shape.has_chart:
            chart = shape.chart
            chart_title = chart.chart_title.text_frame.text if chart.has_title else "No Title"
            text_content = f"[CHART: {chart_title}]"
            
        if text_content:
            side = "RIGHT" if is_on_right else "LEFT"
            out_file.write(f"{indent}- Shape '{shape_name}' ({shape_type}) at {pos_str} [Side: {side}]:\n")
            for line in text_content.split('\n'):
                out_file.write(f"{indent}    {line}\n")
                
    for idx, slide in enumerate(prs.slides, 1):
        out_file.write(f"\nSLIDE {idx}: {slide.shapes.title.text if slide.shapes.title else 'No Title'}\n")
        out_file.write("-" * 50 + "\n")
        
        # Report note/speaker notes
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    out_file.write(f"  [Speaker Notes]:\n    {notes}\n\n")
        except Exception:
            pass
            
        for shape in slide.shapes:
            process_shape(shape, idx)

if __name__ == "__main__":
    out_path = "scratch/pptx_shapes_full_dump.txt"
    with open(out_path, "w", encoding="utf-8") as f:
        dump_all_shapes_and_text("PSAT_Generation_Report_Complete.pptx", f)
    print(f"Dump complete. Output written to {out_path}")

