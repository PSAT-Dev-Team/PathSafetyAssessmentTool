import sys
import json
from pptx import Presentation
from pptx.util import Inches

def analyze_presentation(pptx_path):
    prs = Presentation(pptx_path)
    
    # Slide width and height in inches
    slide_width_in = prs.slide_width / 914400 # 1 inch = 914,400 EMUs
    slide_height_in = prs.slide_height / 914400
    mid_x_in = slide_width_in / 2
    
    print(f"Presentation size: {slide_width_in:.2f} x {slide_height_in:.2f} inches")
    print(f"Midpoint X: {mid_x_in:.2f} inches\n")
    
    analysis_data = []
    
    for idx, slide in enumerate(prs.slides, 1):
        slide_data = {
            "slide_index": idx,
            "title": "",
            "shapes_left_side": [],
            "shapes_right_side": [],
            "native_comments": [],
            "notes": ""
        }
        
        # Try to get slide title
        try:
            if slide.shapes.title:
                slide_data["title"] = slide.shapes.title.text
        except AttributeError:
            pass
            
        # Extract shapes and classify by position (left vs right)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
                
            text = shape.text_frame.text.strip()
            if not text:
                continue
                
            left_in = shape.left / 914400
            top_in = shape.top / 914400
            width_in = shape.width / 914400
            height_in = shape.height / 914400
            
            shape_info = {
                "text": text,
                "box": [left_in, top_in, width_in, height_in],
                "shape_name": shape.name,
                "shape_type": str(shape.shape_type)
            }
            
            # If the box is placed entirely or mostly on the right side of the slide
            # Or if it's explicitly styled as a comment
            is_right = (left_in + width_in / 2) > mid_x_in
            
            if is_right:
                slide_data["shapes_right_side"].append(shape_info)
            else:
                # If it is the title, don't duplicate it in left side if we already captured it,
                # but keep it in slide title or left side for complete context.
                slide_data["shapes_left_side"].append(shape_info)
                
        # Extract native comments if available
        try:
            if hasattr(slide, "comments") and slide.comments:
                for comment in slide.comments:
                    slide_data["native_comments"].append({
                        "author": comment.author,
                        "text": comment.text,
                        "created": str(comment.created) if hasattr(comment, "created") else ""
                    })
        except Exception as e:
            print(f"Error reading native comments on Slide {idx}: {e}")
            
        # Extract notes
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                slide_data["notes"] = slide.notes_slide.notes_text_frame.text.strip()
        except Exception as e:
            pass
            
        analysis_data.append(slide_data)
        
    return analysis_data

if __name__ == "__main__":
    pptx_file = "PSAT_Generation_Report_Complete.pptx"
    try:
        data = analyze_presentation(pptx_file)
        # Write to JSON for inspection
        with open("scratch/pptx_analysis_raw.json", "w", encoding="utf-8") as f:
            json.dump(data, f, indent=2, ensure_ascii=False)
        print("Analysis completed successfully. Output written to scratch/pptx_analysis_raw.json")
    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)
