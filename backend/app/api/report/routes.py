import io
import datetime
from flask import Blueprint, request, send_file, jsonify
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import traceback

from . import bp

@bp.route('/generate-pptx', methods=['POST'])
def generate_pptx():
    try:
        data = request.json
        selected_projects = data.get("selectedProjects", [])
        report_config = data.get("reportConfig", {})
        
        prs = Presentation()
        blank_layout = prs.slide_layouts[5] # Title only
        slide = prs.slides.add_slide(blank_layout)
        
        # Clear the default title placeholder
        for shape in slide.placeholders:
            if shape.is_placeholder:
                shape.text = ""
        
        # --- TITLE SECTION ---
        if report_config.get('showTitle', True):
            if report_config.get('showTitleText', True):
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
                tf = txBox.text_frame
                tf.text = "Path Analysis Executive Summary"
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.size = Pt(28)
                
            if report_config.get('showTitleDescription', True):
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.5))
                tf = txBox.text_frame
                tf.text = f"Analyzed Projects: {', '.join(selected_projects) if selected_projects else 'None'}"
                tf.paragraphs[0].font.size = Pt(14)
                tf.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)

        # --- MAP SECTION ---
        if report_config.get('showMap', True):
            if report_config.get('showMapView', True):
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(4.5), Inches(2.8))
                shape.text = "[ Map Route Visualization Placeholder ]\n" + \
                             "In a full production version, an automated GIS screenshot\n" + \
                             "or static map image would be inserted here."
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
                shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
                shape.text_frame.paragraphs[0].font.size = Pt(12)

        # --- DISTRIBUTION CHARTS ---
        if report_config.get('showCharts', True):
            if report_config.get('showPieChart', True):
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(1.5), Inches(2.1), Inches(2.8))
                shape.text = "[ Pie Chart Placeholder ]"
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
            if report_config.get('showBarChart', True):
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.5), Inches(1.5), Inches(2.1), Inches(2.8))
                shape.text = "[ Bar Chart Placeholder ]"
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(240, 240, 240)

        # --- RISK BANDS SECTION ---
        if report_config.get('showRiskBands', True):
            # We'll use dummy chart data as generating precise pie charts 
            # requires fetching the exact score distributions in a real system.
            y_offset = Inches(4.7)
            
            if report_config.get('showRiskBandsOverall', True):
                txBox = slide.shapes.add_textbox(Inches(0.5), y_offset, Inches(3), Inches(0.5))
                tf = txBox.text_frame
                tf.text = "Overall Risk Level"
                tf.paragraphs[0].font.bold = True
                
                chart_data = ChartData()
                chart_data.categories = ['Low', 'Med', 'High', 'Extreme']
                chart_data.add_series('Risk', (40, 30, 20, 10))
                chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(0.5), y_offset + Inches(0.4), Inches(2), Inches(2), chart_data).chart
                chart.has_legend = False

            if report_config.get('showRiskBandsLegend', True):
                txBox = slide.shapes.add_textbox(Inches(2.6), y_offset, Inches(2), Inches(2))
                tf = txBox.text_frame
                tf.text = "Risk Level Legend:\n- Low Risk\n- Medium Risk\n- High Risk\n- Extreme Risk"
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].font.size = Pt(12)

            if report_config.get('showRiskBandsCrashTypes', True):
                txBox = slide.shapes.add_textbox(Inches(4.5), y_offset, Inches(4), Inches(0.5))
                tf = txBox.text_frame
                tf.text = "Risk by Crash Type"
                tf.paragraphs[0].font.bold = True
                
                # Render the 4 crash types if their toggles are enabled
                cx = Inches(1.2)
                cy = Inches(1.2)
                chart_y = y_offset + Inches(0.6)
                
                if report_config.get('showRiskBandsVB', True):
                    chart_data = ChartData()
                    chart_data.categories = ['L','M','H','E']
                    chart_data.add_series('VB', (20,20,30,30))
                    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(4.5), chart_y, cx, cy, chart_data).chart
                    chart.has_title = True
                    chart.chart_title.text_frame.text = "VB"
                    chart.has_legend = False

                if report_config.get('showRiskBandsBB', True):
                    chart_data = ChartData()
                    chart_data.categories = ['L','M','H','E']
                    chart_data.add_series('BB', (50,20,20,10))
                    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(5.8), chart_y, cx, cy, chart_data).chart
                    chart.has_title = True
                    chart.chart_title.text_frame.text = "BB"
                    chart.has_legend = False

                if report_config.get('showRiskBandsSB', True):
                    chart_data = ChartData()
                    chart_data.categories = ['L','M','H','E']
                    chart_data.add_series('SB', (30,40,20,10))
                    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(7.1), chart_y, cx, cy, chart_data).chart
                    chart.has_title = True
                    chart.chart_title.text_frame.text = "SB"
                    chart.has_legend = False

                if report_config.get('showRiskBandsBP', True):
                    chart_data = ChartData()
                    chart_data.categories = ['L','M','H','E']
                    chart_data.add_series('BP', (60,20,10,10))
                    chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(8.4), chart_y, cx, cy, chart_data).chart
                    chart.has_title = True
                    chart.chart_title.text_frame.text = "BP"
                    chart.has_legend = False

        # Save to a BytesIO object and return
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        return send_file(
            output,
            as_attachment=True,
            download_name="PSAT_Report.pptx",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500


@bp.route('/segment-details', methods=['POST'])
def segment_details():
    """
    For a list of { project, segIndex } references (segIndex is 1-based),
    return each segment's image URL and top 3 contributing risk attributes.
    """
    try:
        import pandas as pd
        from app.services.project_manager import project_manager
        from app.services.cyclerap_scoring import LOOKUP_TABLES

        data = request.json
        segments = data.get("segments", [])

        # Maps LOOKUP_TABLES key → (CSV column name, human display name, safe default)
        ATTR_MAP = {
            'loose_surface':          ("Loose or slippery surface",                           "Loose Surface",           2),
            'tram_rails':             ("Tram or Train Rails",                                 "Tram/Train Rails",        2),
            'delineation':            ("Delineation",                                         "Delineation",             2),
            'light_segregation':      ("Light Segregation",                                   "Light Segregation",       2),
            'facility_width':         ("Facility Width per Direction",                        "Facility Width",          3),
            'flow_direction':         ("Flow Direction",                                       "Flow Direction",          1),
            'width_restriction':      ("Width Restriction",                                   "Width Restriction",       2),
            'adjacent_parking_0_1m':  ("Adjacent Vehicle Parking 0-1m",                      "Adj. Parking 0-1m",       2),
            'adjacent_hazard_0_1m':   ("Adjacent Severe Hazard 0-1m",                        "Adj. Hazard 0-1m",        2),
            'adjacent_parking_1_3m':  ("Adjacent Vehicle Parking 1-3m",                      "Adj. Parking 1-3m",       2),
            'adjacent_hazard_1_3m':   ("Adjacent Severe Hazard 1-3m",                        "Adj. Hazard 1-3m",        2),
            'grade':                  ("Grade",                                               "Grade",                   1),
            'curvature':              ("Curvature",                                           "Curvature",               2),
            'street_lighting':        ("Street Lighting",                                     "Street Lighting",         1),
            'pedestrian_crossing':    ("Pedestrian Crossing",                                 "Pedestrian Crossing",     2),
            'intersecting_facility':  ("Intersecting Bicycle Facility",                       "Intersecting Facility",   2),
            'intersection_crossing':  ("Intersection or Road Crossing",                       "Intersection Crossing",   2),
            'crossing_facility':      ("Crossing Facility",                                   "Crossing Facility",       1),
            'num_lanes_adjacent':     ("Number of lanes – adjacent road",               "Adj. Lane Count",         1),
            'num_lanes_intersecting': ("Number of lanes – intersecting road",           "Intersect. Lanes",        1),
            'property_access':        ("Property Access",                                     "Property Access",         2),
            'pedestrian_flow':        ("Peak pedestrian flow along or across facility",       "Pedestrian Flow",         1),
            'bicycle_flow':           ("Peak bicycle/LV traffic flow",                        "Bicycle Flow",            1),
            'cargo_bikes':            ("Observed proportion of cargo bikes and mopeds",       "Cargo Bikes",             1),
            'bicycle_speed':          ("Bicycle/LV speed – average",                    "Bicycle Speed",           1),
            'speed_differential':     ("Bicycle/LV speed differential",                      "Speed Differential",      1),
            'heavy_vehicle':          ("Heavy vehicle flow",                                  "Heavy Vehicle Flow",      1),
            'line_of_sight':          ("Line of Sight",                                       "Line of Sight",           1),
        }

        pm = project_manager()
        results = []

        for seg_ref in segments:
            project_name = seg_ref.get("project", "")
            seg_index_1based = int(seg_ref.get("segIndex", 1))
            seg_index = seg_index_1based - 1  # convert to 0-based

            try:
                proj = pm.project(project_name)
                ver = proj.latest()
                attrs_df = ver.attributes.df

                if attrs_df is None or seg_index < 0 or seg_index >= len(attrs_df):
                    raise ValueError(f"Segment index {seg_index} out of range")

                row = attrs_df.iloc[seg_index]

                # Try both capitalisation variants of the image reference column
                image_ref = str(row.get("Image reference", "") or row.get("Image Reference", "") or "").strip()
                image_url = f"/api/projects/{project_name}/images/{image_ref}" if image_ref else None

                contributions = []
                for attr_key, (field_name, display_name, default_val) in ATTR_MAP.items():
                    raw_val = row.get(field_name, None)
                    try:
                        val = int(raw_val) if (raw_val is not None and pd.notna(raw_val)) else default_val
                    except (ValueError, TypeError):
                        val = default_val
                    risk = LOOKUP_TABLES.get(attr_key, {}).get(val, {}).get('risk', 1.0)
                    if risk > 1.0:
                        contributions.append({'name': display_name, 'multiplier': round(risk, 2)})

                contributions.sort(key=lambda x: -x['multiplier'])
                results.append({
                    'project': project_name,
                    'segIndex': seg_index_1based,
                    'imageUrl': image_url,
                    'topAttributes': contributions[:3],
                })
            except Exception:
                results.append({
                    'project': project_name,
                    'segIndex': seg_index_1based,
                    'imageUrl': None,
                    'topAttributes': [],
                })

        return jsonify({"ok": True, "details": results})

    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500


@bp.route('/generate-docx', methods=['POST'])
def generate_docx():
    try:
        from docx import Document
        from docx.shared import Pt, Cm, RGBColor as DocxRGB
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from datetime import date

        data = request.json
        selected_projects = data.get("selectedProjects", [])
        elements = data.get("elements", [])
        score_data = data.get("scoreData") or {}
        total_segments = data.get("totalSegments", 0)
        top_risk_rows = data.get("topRiskRows", [])
        oic_name = data.get("oicName", "")
        purpose = data.get("purpose", "")
        report_date = data.get("reportDate", "")
        image_date = data.get("imageDate", "")
        audit_date = data.get("auditDate", "")
        recommendations_text = data.get("recommendations", "")
        score_stats = data.get("scoreStats") or {}
        attribute_frequency = data.get("attributeFrequency") or {}
        project_segment_counts = data.get("projectSegmentCounts") or {}
        project_meta_map = data.get("projectMeta") or {}
        active_filter_names = data.get("activeFilterNames") or []
        all_attribute_rows = data.get("allAttributeRows") or {}
        report_title = data.get("reportTitle") or "Path Analysis Executive Summary"
        project_display_names = data.get("projectDisplayNames") or {}
        section_titles = data.get("sectionTitles") or {}
        map_image_b64 = data.get("mapImageB64")
        deep_dive_image_b64 = data.get("deepDiveImageB64")
        filter_analysis_image_b64 = data.get("filterAnalysisImageB64")

        def disp(name):
            return project_display_names.get(name, name)

        def sec_title(el_id, default_title):
            return section_titles.get(el_id, default_title)

        def insert_b64_image(b64_str, width_cm=15.5):
            import base64
            try:
                img_bytes = base64.b64decode(b64_str)
                doc.add_picture(io.BytesIO(img_bytes), width=Cm(width_cm))
            except Exception:
                pass

        doc = Document()

        # A4 page with comfortable margins
        section = doc.sections[0]
        section.page_width = Cm(21)
        section.page_height = Cm(29.7)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)

        BAND_LABELS = {1: "Low", 2: "Medium", 3: "High", 4: "Extreme"}
        CRASH_TYPE_LABELS = {
            "Overall": "Overall Risk",
            "VB": "Vehicle–Bicycle (VB)",
            "BB": "Bicycle–Bicycle (BB)",
            "SB": "Single-Bicycle (SB)",
            "BP": "Bicycle–Pedestrian (BP)",
        }

        # Sort visible elements by their vertical position on the canvas
        visible_elements = sorted(
            [e for e in elements if e.get("visible", True)],
            key=lambda e: e.get("y", 0)
        )

        for el in visible_elements:
            el_type = el.get("type")

            if el_type == "title":
                h = doc.add_heading(report_title, level=1)
                h.alignment = WD_ALIGN_PARAGRAPH.CENTER
                if selected_projects:
                    p = doc.add_paragraph(f"Projects: {', '.join(selected_projects)}")
                    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
                if oic_name:
                    p_oic = doc.add_paragraph()
                    p_oic.add_run("OIC In-charge: ").bold = True
                    p_oic.add_run(oic_name)
                if purpose:
                    p_pur = doc.add_paragraph()
                    p_pur.add_run("Purpose: ").bold = True
                    p_pur.add_run(purpose)

                def fmt_iso_date(iso_str):
                    if not iso_str:
                        return ""
                    try:
                        from datetime import datetime as _dt
                        return _dt.fromisoformat(iso_str).strftime("%d %B %Y")
                    except Exception:
                        return iso_str

                if report_date:
                    p_rd = doc.add_paragraph()
                    p_rd.add_run("Report Date: ").bold = True
                    p_rd.add_run(fmt_iso_date(report_date))
                if image_date:
                    p_id = doc.add_paragraph()
                    p_id.add_run("Image Date: ").bold = True
                    p_id.add_run(fmt_iso_date(image_date))
                if audit_date:
                    p_ad = doc.add_paragraph()
                    p_ad.add_run("Audit Date: ").bold = True
                    p_ad.add_run(fmt_iso_date(audit_date))

                p_date = doc.add_paragraph(f"Generated: {date.today().strftime('%d %B %Y')}")
                p_date.alignment = WD_ALIGN_PARAGRAPH.CENTER
                doc.add_paragraph()

            elif el_type == "summary":
                doc.add_heading(sec_title(el.get("id",""), "Summary"), level=2)
                p1 = doc.add_paragraph()
                p1.add_run("Projects Analyzed: ").bold = True
                p1.add_run(str(len(selected_projects)))
                p2 = doc.add_paragraph()
                p2.add_run("Total Segments: ").bold = True
                p2.add_run(str(total_segments))
                doc.add_paragraph()

            elif el_type == "riskBands":
                doc.add_heading(sec_title(el.get("id",""), "Risk Band Distribution"), level=2)
                if score_data:
                    for crash_key, crash_label in CRASH_TYPE_LABELS.items():
                        dist = score_data.get(crash_key, {})
                        if not dist:
                            continue
                        # Accept both string and int keys
                        total = sum(int(v) for v in dist.values())
                        if total == 0:
                            continue

                        doc.add_heading(crash_label, level=3)
                        table = doc.add_table(rows=1, cols=3)
                        table.style = "Table Grid"
                        hdr = table.rows[0].cells
                        hdr[0].text = "Risk Band"
                        hdr[1].text = "Segments"
                        hdr[2].text = "Percentage"
                        for band_num in [1, 2, 3, 4]:
                            count = int(dist.get(band_num, dist.get(str(band_num), 0)))
                            pct = (count / total * 100) if total > 0 else 0
                            row = table.add_row().cells
                            row[0].text = BAND_LABELS.get(band_num, str(band_num))
                            row[1].text = str(count)
                            row[2].text = f"{pct:.1f}%"
                        doc.add_paragraph()
                else:
                    doc.add_paragraph("No score data available.")
                    doc.add_paragraph()

            elif el_type == "map":
                doc.add_heading(sec_title(el.get("id",""), "Map View"), level=2)
                if map_image_b64:
                    insert_b64_image(map_image_b64)
                else:
                    p = doc.add_paragraph("[Map image not available — export again from the Report Builder]")
                    p.runs[0].italic = True
                    p.runs[0].font.color.rgb = DocxRGB(0x88, 0x88, 0x88)
                doc.add_paragraph()

            elif el_type == "topRisk":
                doc.add_heading(sec_title(el.get("id",""), "Top Risk Stretches"), level=2)
                if top_risk_rows:
                    treatment_names_map = data.get("treatmentNames", {})
                    table = doc.add_table(rows=1, cols=9)
                    table.style = "Table Grid"
                    hdr = table.rows[0].cells
                    for i, h in enumerate(["#", "Project", "Seg", "Score", "VB", "BB", "SB", "BP", "Applied Treatments"]):
                        hdr[i].text = h
                    BAND_LABELS_SHORT = {1: "Low", 2: "Med", 3: "High", 4: "Extreme"}
                    for rank, row in enumerate(top_risk_rows, start=1):
                        cells = table.add_row().cells
                        cells[0].text = str(rank)
                        cells[1].text = disp(str(row.get("_project", "")))
                        cells[2].text = str(row.get("_segIndex", ""))
                        max_score = row.get("_maxScore", 0)
                        cells[3].text = f"{float(max_score):.1f}" if max_score else "—"
                        cells[4].text = BAND_LABELS_SHORT.get(row.get("VB Band", 0), "—")
                        cells[5].text = BAND_LABELS_SHORT.get(row.get("BB Band", 0), "—")
                        cells[6].text = BAND_LABELS_SHORT.get(row.get("SB Band", 0), "—")
                        cells[7].text = BAND_LABELS_SHORT.get(row.get("BP Band", 0), "—")
                        applied = row.get("_treatments", [])
                        if applied:
                            treatment_lines = []
                            for t_id in applied:
                                t_name = treatment_names_map.get(
                                    str(t_id),
                                    treatment_names_map.get(int(t_id) if str(t_id).isdigit() else t_id, f"Treatment {t_id}")
                                )
                                treatment_lines.append(f"{t_id}. {t_name}")
                            cells[8].text = "\n".join(treatment_lines)
                        else:
                            cells[8].text = "None"
                    doc.add_paragraph()
                else:
                    doc.add_paragraph("No segment score data available.")
                    doc.add_paragraph()

            elif el_type == "treatmentSummary":
                treatment_summaries = data.get("treatmentSummaries", [])
                treatment_names = data.get("treatmentNames", {})

                doc.add_heading(sec_title(el.get("id",""), "Treatment Summary"), level=2)

                if not treatment_summaries:
                    doc.add_paragraph("No treatment data available.")
                    doc.add_paragraph()
                else:
                    for summary in treatment_summaries:
                        project_name = summary.get("project", "Unknown")
                        treated_segs = summary.get("treatedSegments", 0)
                        treatment_counts = summary.get("treatmentCounts", {})

                        doc.add_heading(disp(project_name), level=3)

                        p = doc.add_paragraph()
                        p.add_run("Treated Segments: ").bold = True
                        p.add_run(str(treated_segs))

                        if treatment_counts:
                            tbl = doc.add_table(rows=1, cols=3)
                            tbl.style = "Table Grid"
                            hdr = tbl.rows[0].cells
                            hdr[0].text = "ID"
                            hdr[1].text = "Treatment"
                            hdr[2].text = "Segments Applied"

                            sorted_treatments = sorted(
                                treatment_counts.items(),
                                key=lambda x: -int(x[1])
                            )
                            for t_id_str, count in sorted_treatments:
                                t_name = treatment_names.get(
                                    t_id_str,
                                    treatment_names.get(int(t_id_str) if str(t_id_str).isdigit() else t_id_str, f"Treatment {t_id_str}")
                                )
                                row_cells = tbl.add_row().cells
                                row_cells[0].text = str(t_id_str)
                                row_cells[1].text = str(t_name)
                                row_cells[2].text = str(count)
                            doc.add_paragraph()
                        else:
                            doc.add_paragraph("No treatments applied yet.")
                            doc.add_paragraph()

            elif el_type == "projectDetails":
                doc.add_heading(sec_title(el.get("id",""), "Project Details"), level=2)
                if selected_projects:
                    for proj_name in selected_projects:
                        doc.add_heading(disp(proj_name), level=3)
                        meta = project_meta_map.get(proj_name, {})
                        seg_count = project_segment_counts.get(proj_name, 0)

                        def add_detail_row(label, value):
                            p = doc.add_paragraph()
                            p.add_run(f"{label}: ").bold = True
                            p.add_run(str(value))

                        add_detail_row("Segments", seg_count)
                        length_km = meta.get("lengthKm")
                        if length_km is not None:
                            add_detail_row("Length", f"{float(length_km):.1f} km")

                        def fmt_date(iso_str):
                            if not iso_str:
                                return "—"
                            try:
                                from datetime import datetime
                                dt = datetime.fromisoformat(str(iso_str).replace("Z", "+00:00"))
                                return dt.strftime("%-d %b %Y") if hasattr(dt, 'strftime') else str(iso_str)
                            except Exception:
                                return str(iso_str)

                        add_detail_row("Survey",   fmt_date(meta.get("dateCreated")))
                        add_detail_row("Analysis", fmt_date(meta.get("lastUpdated")))
                    doc.add_paragraph()
                else:
                    doc.add_paragraph("No project data available.")
                    doc.add_paragraph()

            elif el_type == "riskStats":
                doc.add_heading(sec_title(el.get("id",""), "Risk Score Statistics"), level=2)
                if score_stats:
                    tbl = doc.add_table(rows=1, cols=4)
                    tbl.style = "Table Grid"
                    hdr = tbl.rows[0].cells
                    for i, h in enumerate(["Crash Type", "Min Score", "Max Score", "Avg Score"]):
                        hdr[i].text = h
                    crash_labels = {"Overall": "Overall Risk", "VB": "Vehicle–Bicycle",
                                    "BB": "Bicycle–Bicycle", "SB": "Single-Bicycle", "BP": "Bicycle–Pedestrian"}
                    for ct_key, ct_label in crash_labels.items():
                        stats = score_stats.get(ct_key, {})
                        row_cells = tbl.add_row().cells
                        row_cells[0].text = ct_label
                        row_cells[1].text = str(stats.get("min", "—"))
                        row_cells[2].text = str(stats.get("max", "—"))
                        row_cells[3].text = str(stats.get("avg", "—"))
                    doc.add_paragraph()
                else:
                    doc.add_paragraph("No score statistics available.")
                    doc.add_paragraph()

            elif el_type == "topAttributes":
                doc.add_heading(sec_title(el.get("id",""), "Top Risk Factors"), level=2)
                doc.add_paragraph("Most frequently contributing risk attributes among top-risk segments:")
                if attribute_frequency:
                    tbl = doc.add_table(rows=1, cols=2)
                    tbl.style = "Table Grid"
                    hdr = tbl.rows[0].cells
                    hdr[0].text = "Risk Factor"
                    hdr[1].text = "Segments Affected"
                    for attr_name, count in sorted(attribute_frequency.items(), key=lambda x: -int(x[1])):
                        row_cells = tbl.add_row().cells
                        row_cells[0].text = str(attr_name)
                        row_cells[1].text = str(count)
                    doc.add_paragraph()
                else:
                    doc.add_paragraph("No attribute frequency data available.")
                    doc.add_paragraph()

            elif el_type == "recommendations":
                doc.add_heading(sec_title(el.get("id",""), "Recommendations"), level=2)
                if recommendations_text and recommendations_text.strip():
                    for line in recommendations_text.strip().split("\n"):
                        doc.add_paragraph(line if line.strip() else "")
                else:
                    p = doc.add_paragraph("[No recommendations entered]")
                    p.runs[0].italic = True
                    p.runs[0].font.color.rgb = DocxRGB(0xAA, 0xAA, 0xAA)
                doc.add_paragraph()

            elif el_type == "methodology":
                doc.add_heading(sec_title(el.get("id",""), "Methodology"), level=2)
                doc.add_heading("CycleRAP v2 — Cycling Road Assessment Programme", level=3)
                methodology_body = (
                    "This report uses the CycleRAP (Cycling Road Assessment Programme) methodology to assess "
                    "the safety of cycling infrastructure. Each segment is evaluated against a set of risk "
                    "attributes covering facility design, surface quality, hazards, intersections, and usage "
                    "patterns. A risk multiplier is computed for each attribute based on its coded value, and "
                    "the combined score determines the segment's risk band (Low / Medium / High / Extreme) for "
                    "four crash types: Vehicle–Bicycle (VB), Bicycle–Bicycle (BB), Single-Bicycle (SB), and "
                    "Bicycle–Pedestrian (BP). Higher scores and bands indicate greater risk exposure and a "
                    "greater need for intervention."
                )
                doc.add_paragraph(methodology_body)
                doc.add_paragraph()

            elif el_type == "segmentGallery":
                doc.add_heading(sec_title(el.get("id",""), "Segment Image Gallery"), level=2)
                doc.add_paragraph(
                    "[Segment images are shown in the PDF export. In the Word document, please insert "
                    "segment images manually from the Path Analysis page or Image Gallery export.]"
                ).runs[0].italic = True
                doc.add_paragraph()

            elif el_type == "deepDive":
                doc.add_heading(sec_title(el.get("id",""), "Deep-Dive Risk Analytics"), level=2)

                if deep_dive_image_b64:
                    insert_b64_image(deep_dive_image_b64)
                    doc.add_paragraph()

                # Overall risk distribution table
                doc.add_heading("Overall Risk Distribution", level=3)
                overall_dist = (score_data.get("Overall") or {}) if score_data else {}
                BAND_NAMES = {1: "Low", 2: "Medium", 3: "High", 4: "Extreme"}
                if overall_dist:
                    total_segs = sum(int(v) for v in overall_dist.values())
                    tbl = doc.add_table(rows=1, cols=3)
                    tbl.style = "Table Grid"
                    hdr = tbl.rows[0].cells
                    hdr[0].text = "Risk Band"
                    hdr[1].text = "Segments"
                    hdr[2].text = "Percentage"
                    for band_num in [1, 2, 3, 4]:
                        count = int(overall_dist.get(band_num, overall_dist.get(str(band_num), 0)))
                        pct = f"{count / total_segs * 100:.1f}%" if total_segs > 0 else "—"
                        row_cells = tbl.add_row().cells
                        row_cells[0].text = BAND_NAMES.get(band_num, str(band_num))
                        row_cells[1].text = str(count)
                        row_cells[2].text = pct
                else:
                    doc.add_paragraph("No score data available.").runs[0].italic = True
                doc.add_paragraph()

                # Top contributing attributes table
                doc.add_heading("Top Contributing Attributes Across Project", level=3)
                if attribute_frequency:
                    tbl2 = doc.add_table(rows=1, cols=2)
                    tbl2.style = "Table Grid"
                    hdr2 = tbl2.rows[0].cells
                    hdr2[0].text = "Attribute"
                    hdr2[1].text = "Segments Affected"
                    for attr_name, count in sorted(attribute_frequency.items(), key=lambda x: -int(x[1])):
                        row_cells = tbl2.add_row().cells
                        row_cells[0].text = str(attr_name)
                        row_cells[1].text = str(count)
                else:
                    doc.add_paragraph("No attribute frequency data available.").runs[0].italic = True
                doc.add_paragraph()

            elif el_type == "filterAnalysis":
                doc.add_heading(sec_title(el.get("id",""), "Filter Analysis"), level=2)

                if filter_analysis_image_b64:
                    insert_b64_image(filter_analysis_image_b64)
                    doc.add_paragraph()

                if not active_filter_names:
                    doc.add_paragraph("No active filters were applied in Path Analysis.").runs[0].italic = True
                    doc.add_paragraph()
                else:
                    # Combine all attribute rows across projects
                    all_rows_combined = []
                    for proj_rows in all_attribute_rows.values():
                        all_rows_combined.extend(proj_rows)

                    for filter_name in active_filter_names:
                        doc.add_heading(filter_name, level=3)
                        value_counts = {}
                        for row_data in all_rows_combined:
                            val = row_data.get(filter_name)
                            if val is not None and str(val).strip() != "":
                                key = str(val)
                                value_counts[key] = value_counts.get(key, 0) + 1

                        if not value_counts:
                            doc.add_paragraph("No data for this attribute.").runs[0].italic = True
                        else:
                            tbl_f = doc.add_table(rows=1, cols=2)
                            tbl_f.style = "Table Grid"
                            hdr_f = tbl_f.rows[0].cells
                            hdr_f[0].text = "Value"
                            hdr_f[1].text = "Segments"
                            for val_key in sorted(value_counts.keys(), key=lambda x: (int(x) if x.lstrip("-").isdigit() else x)):
                                r_cells = tbl_f.add_row().cells
                                r_cells[0].text = val_key
                                r_cells[1].text = str(value_counts[val_key])
                    doc.add_paragraph()

        output = io.BytesIO()
        doc.save(output)
        output.seek(0)

        return send_file(
            output,
            as_attachment=True,
            download_name="PSAT_Report.docx",
            mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )

    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500
