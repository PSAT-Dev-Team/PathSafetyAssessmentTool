from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData, CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import datetime

def create_presentation():
    prs = Presentation()
    blank_layout = prs.slide_layouts[5] # Title only

    # --- Title Slide ---
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "PSAT Generation Report"
    slide.placeholders[1].text = "Expanded Draft Concepts, Layout Ideas & UI Flow\n" + datetime.datetime.now().strftime("%Y-%m-%d")

    # --- NEW IDEA: UI Integration & User Flow ---
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.title.text = "UI Mockup: How Users Will Generate The Report"
    
    # Draw a mock "App Interface"
    # Main App Window
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(8.5), Inches(5)).fill.solid()
    # Sidebar
    sidebar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(2), Inches(5))
    sidebar.fill.solid()
    sidebar.fill.fore_color.rgb = RGBColor(220, 220, 220)
    
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(1.8), Inches(4))
    tf = txBox.text_frame
    tf.text = "PSAT Menu\n\n- Dashboard\n- GIS Layers\n- Projects"
    
    # Print / Export Button in Sidebar
    btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.5), Inches(1.8), Inches(0.6))
    btn.fill.solid()
    btn.fill.fore_color.rgb = RGBColor(0, 112, 192) # Blue button
    btn_tf = btn.text_frame
    btn_tf.text = "📄 Generate Report"
    btn_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    btn_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    btn_tf.paragraphs[0].font.bold = True
    
    # Main Content Area
    txBox = slide.shapes.add_textbox(Inches(2.7), Inches(1.6), Inches(6), Inches(2))
    tf = txBox.text_frame
    tf.text = "Project: Central District Cycle Path\nOverall Status: Completed\n\n(Map and Project Details loaded here...)"
    
    # Callout Arrow explaining the button
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(2.6), Inches(5.6), Inches(1.5), Inches(0.4))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(255, 192, 0)
    
    txBox = slide.shapes.add_textbox(Inches(4.3), Inches(5.4), Inches(4), Inches(1.5))
    tf = txBox.text_frame
    tf.text = "User Journey:\n1. User views completed project.\n2. Clicks 'Generate Report' in sidebar.\n3. System complies the Word Doc/PPTX."
    tf.paragraphs[0].font.bold = True

    # --- Idea 1: Project Overview (Dashboard Style) ---
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.title.text = "Concept 1: Project Overview & Map Location"
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(4))
    tf = txBox.text_frame
    tf.text = "Project Details"
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.size = Pt(24)
    details = ["Project Name: Central District", "Segments: 45", "Length: 12.5 km", "Survey: Oct 12, 2025", "Analysis: Oct 15, 2025"]
    for detail in details: 
        p = tf.add_paragraph()
        p.text = detail
        p.font.size = Pt(18)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(1.5), Inches(4.5), Inches(5))
    shape.text = "[ Map Location Image Placeholder ]"

    # --- Idea 2: Combined Executive Summary (One-Pager) ---
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.title.text = "Concept 2: Combined Executive Summary (One-Pager)"
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(3), Inches(2))
    tf = txBox.text_frame
    tf.text = "Central District Project\nSegments: 45 | 12.5 km\nAnalyzed: Oct 15, 2025"
    tf.paragraphs[0].font.bold = True
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(1.5), Inches(5.5), Inches(2.5))
    shape.text = "[ Wide Map Image Overview ]"
    
    chart_data = ChartData()
    chart_data.categories = ['High', 'Med', 'Low']
    chart_data.add_series('Risk', (15, 35, 50))
    chart1 = slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(0.5), Inches(4.5), Inches(4), Inches(2.5), chart_data).chart
    chart1.has_title = True
    chart1.chart_title.text_frame.text = "Bicycle-Bicycle"

    chart_data2 = ChartData()
    chart_data2.categories = ['High', 'Med', 'Low']
    chart_data2.add_series('Risk', (25, 45, 30))
    chart2 = slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(5), Inches(4.5), Inches(4), Inches(2.5), chart_data2).chart
    chart2.has_title = True
    chart2.chart_title.text_frame.text = "Bicycle-Pedestrian"

    # --- Idea 3: Deep-Dive Risk Analytics ---
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.title.text = "Concept 3: Deep-Dive Risk Analytics"
    
    chart1 = slide.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(0.5), Inches(1.5), Inches(4), Inches(4), chart_data).chart
    chart1.has_title = True
    chart1.chart_title.text_frame.text = "Overall Risk Distribution"
    
    bar_data = CategoryChartData()
    bar_data.categories = ['Narrow Width', 'Obstacle', 'Slope', 'Surface', 'Visibility']
    bar_data.add_series('Frequency', (24, 18, 12, 8, 5))
    chart3 = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(5), Inches(1.5), Inches(4.5), Inches(4), bar_data).chart
    chart3.has_title = True
    chart3.chart_title.text_frame.text = "Top Contributing Attributes Across Project"

    # --- Idea 4: Top 10 High Risk Stretches (Detailed List) ---
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.title.text = "Concept 4: Top Risk Stretches (List View)"
    y_pos = Inches(1.5)
    for i in range(1, 4):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y_pos, Inches(2), Inches(1.5))
        shape.text = f"Image #{i}"
        txBox = slide.shapes.add_textbox(Inches(2.7), y_pos, Inches(6), Inches(1.5))
        tf = txBox.text_frame
        tf.text = f"#{i} High Risk Stretch - Score: 4.5"
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
        p = tf.add_paragraph()
        p.text = "Attributes: 1. Narrow Width | 2. Obstacle | 3. Slope"
        y_pos += Inches(1.8)
        
    # --- Idea 5: Top 10 High Risk Stretches (Grid/Cards) ---
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.title.text = "Concept 5: Top Risk Stretches (Grid View)"
    positions = [(Inches(0.5), Inches(1.5)), (Inches(3.8), Inches(1.5)), (Inches(7.1), Inches(1.5)),
                 (Inches(0.5), Inches(4.5)), (Inches(3.8), Inches(4.5)), (Inches(7.1), Inches(4.5))]
    for i, (x, y) in enumerate(positions, 1):
        slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3), Inches(2.8))
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x+Inches(0.2), y+Inches(0.2), Inches(2.6), Inches(1))
        shape.text = f"Img {i}"
        txBox = slide.shapes.add_textbox(x, y+Inches(1.2), Inches(3), Inches(1.5))
        tf = txBox.text_frame
        tf.text = f"#{i} Score: 4.5"
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
        p = tf.add_paragraph()
        p.text = "- Narrow Width\n- Obstacle\n- Slope: 8%"
        p.font.size = Pt(12)

    # --- Idea 6: Top 10 High Risk Stretches (Tabular Report View) ---
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.title.text = "Concept 6: Top Risk Stretches (Tabular View)"
    
    rows, cols = 6, 4
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(5)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    table.columns[0].width = Inches(1.0)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(1.5)
    table.columns[3].width = Inches(4.0)
    
    headers = ["Rank", "Image", "Risk Score", "Top 3 Contributing Attributes"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        
    for r in range(1, 6):
        table.cell(r, 0).text = f"#{r}"
        table.cell(r, 1).text = "[ Image Placeholder ]"
        table.cell(r, 2).text = "High (4.5)"
        table.cell(r, 2).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
        table.cell(r, 3).text = "1. Narrow Width\n2. Obstacle\n3. Slope"

    # --- Idea 7: Infographic Style Breakdown ---
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.title.text = "Concept 7: Infographic / Data Callout Style"
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(2.5), Inches(2))
    tf = txBox.text_frame
    tf.text = "45"
    tf.paragraphs[0].font.size = Pt(60)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "Segments Analyzed"
    p.alignment = PP_ALIGN.CENTER

    txBox = slide.shapes.add_textbox(Inches(3.5), Inches(1.5), Inches(2.5), Inches(2))
    tf = txBox.text_frame
    tf.text = "12.5 km"
    tf.paragraphs[0].font.size = Pt(60)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "Total Facility Length"
    p.alignment = PP_ALIGN.CENTER
    
    txBox = slide.shapes.add_textbox(Inches(6.5), Inches(1.5), Inches(2.5), Inches(2))
    tf = txBox.text_frame
    tf.text = "34%"
    tf.paragraphs[0].font.size = Pt(60)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.add_paragraph()
    p.text = "High Risk Corridors"
    p.alignment = PP_ALIGN.CENTER
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(4), Inches(8), Inches(2.5))
    shape.text = "[ Integrated Map with Heatmap Overlay Placeholder ]"

    # Save presentation
    output_path = "PSAT_Generation_Report_Complete.pptx"
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == '__main__':
    create_presentation()
